"""
Office Tools: docx and pptx reading and writing.
"""

import os
from pathlib import Path
from core.skills import SkillManager


def _ensure_dependencies():
    """Dynamically install python-docx and python-pptx if missing."""
    try:
        import docx
        import pptx
    except ImportError:
        import subprocess
        import sys
        print("  [OFFICE] Installing missing dependencies: python-docx, python-pptx...")
        subprocess.run([sys.executable, "-m", "pip", "install", "python-docx", "python-pptx", "-q"], check=False)


def register(manager: SkillManager) -> None:
    """Register Office-related skills."""
    _ensure_dependencies()

    @manager.skill(
        name="read_docx",
        description="读取 Word (.docx) 文件的纯文本内容。",
        parameters={
            "properties": {
                "path": {"type": "string", "description": "文件路径"}
            },
            "required": ["path"]
        },
        category="office"
    )
    def read_docx(path: str) -> str:
        try:
            from docx import Document
            p = Path(path)
            if not p.exists(): return f"错误: 文件不存在 {path}"
            doc = Document(str(p))
            full_text = [para.text for para in doc.paragraphs]
            return "\n".join(full_text)
        except Exception as e:
            return f"读取失败: {e}"

    @manager.skill(
        name="create_docx",
        description="创建一个包含指定内容的 Word (.docx) 文件。",
        parameters={
            "properties": {
                "path": {"type": "string", "description": "保存路径"},
                "title": {"type": "string", "description": "文档标题"},
                "paragraphs": {"type": "array", "items": {"type": "string"}, "description": "正文段落列表"}
            },
            "required": ["path", "paragraphs"]
        },
        category="office"
    )
    def create_docx(path: str, paragraphs: list, title: str = None) -> str:
        try:
            from docx import Document
            doc = Document()
            if title: doc.add_heading(title, 0)
            for p_text in paragraphs:
                doc.add_paragraph(p_text)
            doc.save(path)
            return f"[OK] Word 文档已保存到: {path}"
        except Exception as e:
            return f"创建失败: {e}"

    @manager.skill(
        name="read_pptx",
        description="读取 PowerPoint (.pptx) 演示文稿的所有幻灯片文本。",
        parameters={
            "properties": {
                "path": {"type": "string", "description": "文件路径"}
            },
            "required": ["path"]
        },
        category="office"
    )
    def read_pptx(path: str) -> str:
        try:
            from pptx import Presentation
            p = Path(path)
            if not p.exists(): return f"错误: 文件不存在 {path}"
            prs = Presentation(str(p))
            text_runs = []
            for i, slide in enumerate(prs.slides):
                text_runs.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            return f"读取失败: {e}"

    @manager.skill(
        name="create_pptx",
        description="从幻灯片内容列表创建一个 PowerPoint (.pptx) 文件。",
        parameters={
            "properties": {
                "path": {"type": "string", "description": "保存路径"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"}
                        }
                    },
                    "description": "幻灯片列表"
                }
            },
            "required": ["path", "slides"]
        },
        category="office"
    )
    def create_pptx(path: str, slides: list) -> str:
        try:
            from pptx import Presentation
            prs = Presentation()
            for s_data in slides:
                slide_layout = prs.slide_layouts[1] # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = s_data.get("title", "")
                content.text = s_data.get("content", "")
            prs.save(path)
            return f"[OK] PPT 文档已保存到: {path}"
        except Exception as e:
            return f"创建失败: {e}"
